#!/usr/bin/env python3
import argparse
import base64
import json
import os
import ssl
import subprocess
import tempfile
from functools import partial
from http.server import SimpleHTTPRequestHandler, ThreadingHTTPServer

WIDESCREEN_WIDTH_EMU = 12192000
WIDESCREEN_HEIGHT_EMU = 6858000


def applescript_quote(value: str) -> str:
    return value.replace("\\", "\\\\").replace('"', '\\"')


def run_osascript(script: str) -> str:
    result = subprocess.run(
        ["osascript", "-"],
        input=script,
        text=True,
        capture_output=True,
        check=False,
    )
    if result.returncode != 0:
        raise RuntimeError(result.stderr.strip() or "osascript failed")
    return result.stdout.strip()


def choose_save_path(suggested_filename: str) -> str:
    script = f'''
set chosenAlias to choose file name with prompt "Save Jolify export" default name "{applescript_quote(suggested_filename)}"
POSIX path of chosenAlias
'''.strip()
    return run_osascript(script)


def save_base64_file(base64_file: str, suggested_filename: str) -> str:
    output_path = choose_save_path(suggested_filename)
    with open(output_path, "wb") as handle:
        handle.write(base64.b64decode(base64_file))
    return output_path


def open_in_powerpoint(path: str) -> None:
    subprocess.run(["open", "-a", "Microsoft PowerPoint", path], check=True)


def create_outlook_draft(path: str, subject: str) -> None:
    script = f'''
set attachmentFile to POSIX file "{applescript_quote(path)}"
tell application "Microsoft Outlook"
  activate
  set draftMessage to make new outgoing message with properties {{subject:"{applescript_quote(subject)}"}}
  make new attachment at draftMessage with properties {{file:attachmentFile}}
  open draftMessage
end tell
'''.strip()
    run_osascript(script)


def create_presentation_from_pictures(images: list[dict], suggested_filename: str) -> str:
    if not images:
        raise RuntimeError("At least one image is required.")

    try:
        from PIL import Image
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError as error:
        raise RuntimeError(
            "Picture deck creation requires Jolify local Python dependencies. "
            "Please rerun install-local.sh to repair the local runtime."
        ) from error

    with tempfile.TemporaryDirectory(prefix="jolify-pictures-") as temp_dir:
        prs = Presentation()
        prs.slide_width = WIDESCREEN_WIDTH_EMU
        prs.slide_height = WIDESCREEN_HEIGHT_EMU
        blank_layout = prs.slide_layouts[6]

        for image in images:
            filename = image.get("filename") or "image.png"
            base64_image = image.get("base64Image")
            if not isinstance(base64_image, str) or not base64_image:
                raise RuntimeError("Each image must include base64Image.")

            image_path = os.path.join(temp_dir, os.path.basename(filename))
            with open(image_path, "wb") as handle:
                handle.write(base64.b64decode(base64_image))

            with Image.open(image_path) as source:
                image_width_px, image_height_px = source.size

            slide = prs.slides.add_slide(blank_layout)

            scale = min(
                WIDESCREEN_WIDTH_EMU / image_width_px,
                WIDESCREEN_HEIGHT_EMU / image_height_px,
            )
            fitted_width = int(image_width_px * scale)
            fitted_height = int(image_height_px * scale)
            left = int((WIDESCREEN_WIDTH_EMU - fitted_width) / 2)
            top = int((WIDESCREEN_HEIGHT_EMU - fitted_height) / 2)

            slide.shapes.add_picture(
                image_path,
                Emu(left),
                Emu(top),
                width=Emu(fitted_width),
                height=Emu(fitted_height),
            )

        output_path = choose_save_path(suggested_filename)
        prs.save(output_path)
        return output_path


class JolifyHandler(SimpleHTTPRequestHandler):
    def end_headers(self):
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Headers", "Content-Type")
        self.send_header("Access-Control-Allow-Methods", "GET, POST, OPTIONS")
        super().end_headers()

    def send_json(self, status: int, payload: dict):
        body = json.dumps(payload).encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)

    def do_OPTIONS(self):
        self.send_response(204)
        self.end_headers()

    def do_GET(self):
        if self.path in ("/healthz", "/healthz/"):
            body = b"ok\n"
            self.send_response(200)
            self.send_header("Content-Type", "text/plain; charset=utf-8")
            self.send_header("Content-Length", str(len(body)))
            self.end_headers()
            self.wfile.write(body)
            return

        super().do_GET()

    def do_POST(self):
        if self.path not in ("/native/save-file", "/native/create-outlook-draft", "/native/create-presentation-from-pictures"):
          self.send_json(404, {"error": "Unknown endpoint."})
          return

        try:
            content_length = int(self.headers.get("Content-Length", "0"))
            body = self.rfile.read(content_length)
            payload = json.loads(body.decode("utf-8") or "{}")

            if self.path == "/native/save-file":
                base64_file = payload.get("base64File")
                suggested_filename = payload.get("suggestedFilename") or "Jolify Export.pptx"
                open_after_save = bool(payload.get("openInPowerPoint"))

                if not isinstance(base64_file, str) or not base64_file:
                    self.send_json(400, {"error": "base64File is required."})
                    return

                saved_path = save_base64_file(base64_file, suggested_filename)
                if open_after_save:
                    open_in_powerpoint(saved_path)

                self.send_json(200, {"savedPath": saved_path})
                return

            if self.path == "/native/create-presentation-from-pictures":
                images = payload.get("images")
                suggested_filename = payload.get("suggestedFilename") or "Pictures Presentation.pptx"

                if not isinstance(images, list) or len(images) == 0:
                    self.send_json(400, {"error": "images is required."})
                    return

                saved_path = create_presentation_from_pictures(images, suggested_filename)
                open_in_powerpoint(saved_path)
                self.send_json(200, {"savedPath": saved_path})
                return

            attachment_path = payload.get("attachmentPath")
            subject = payload.get("subject") or "Jolify export"
            if not isinstance(attachment_path, str) or not attachment_path:
                self.send_json(400, {"error": "attachmentPath is required."})
                return

            if not os.path.exists(attachment_path):
                self.send_json(400, {"error": "Attachment file does not exist."})
                return

            create_outlook_draft(attachment_path, subject)
            self.send_json(200, {"ok": True})
        except RuntimeError as error:
            self.send_json(500, {"error": str(error)})
        except subprocess.CalledProcessError as error:
            self.send_json(500, {"error": error.stderr or error.stdout or str(error)})
        except Exception as error:
            self.send_json(500, {"error": str(error)})

    def log_message(self, format, *args):
        print("%s - - [%s] %s" % (self.address_string(), self.log_date_time_string(), format % args), flush=True)


def main():
    parser = argparse.ArgumentParser(description="Serve Jolify locally over HTTPS.")
    parser.add_argument("--root", required=True, help="Directory to serve")
    parser.add_argument("--host", default="127.0.0.1", help="Host to bind")
    parser.add_argument("--port", type=int, default=38443, help="Port to bind")
    parser.add_argument("--cert", required=True, help="TLS certificate path")
    parser.add_argument("--key", required=True, help="TLS private key path")
    args = parser.parse_args()

    handler = partial(JolifyHandler, directory=args.root)
    httpd = ThreadingHTTPServer((args.host, args.port), handler)
    context = ssl.SSLContext(ssl.PROTOCOL_TLS_SERVER)
    context.load_cert_chain(certfile=args.cert, keyfile=args.key)
    httpd.socket = context.wrap_socket(httpd.socket, server_side=True)

    print(f"Serving Jolify from {args.root} on https://{args.host}:{args.port}", flush=True)
    httpd.serve_forever()


if __name__ == "__main__":
    main()
